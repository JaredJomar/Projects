from __future__ import annotations

import datetime
import os
from pathlib import Path
from typing import Iterable, List, Sequence

import fitz  # PyMuPDF
import pandas as pd
import pdfplumber
import PyPDF2
import tabula
from pdf2docx import Converter
from pptx import Presentation
from pptx.util import Inches
from tabula.errors import JavaNotFoundError

__all__ = [
    "PdfOperationError",
    "parse_page_ranges",
    "merge_pdfs",
    "extract_pages",
    "split_pdf",
    "rotate_pdf",
    "compress_pdf",
    "convert_pdf",
    "get_pdf_page_count",
]


class PdfOperationError(Exception):
    """Raised when a PDF operation fails."""


def _ensure_output_dir(path: os.PathLike[str] | str) -> Path:
    directory = Path(path)
    directory.mkdir(parents=True, exist_ok=True)
    return directory


def parse_page_ranges(range_string: str) -> List[int]:
    """Parse strings like "1,3,5-7" into a sorted list of unique page numbers."""
    if not range_string.strip():
        return []

    pages: set[int] = set()
    for raw_part in range_string.split(','):
        part = raw_part.strip()
        if not part:
            continue
        if '-' in part:
            try:
                start_str, end_str = part.split('-', maxsplit=1)
                start = int(start_str.strip())
                end = int(end_str.strip())
            except ValueError as exc:
                raise PdfOperationError(f"Invalid range: {part}") from exc
            if start > end:
                raise PdfOperationError(f"Invalid range: {part} (start > end)")
            pages.update(range(start, end + 1))
        else:
            try:
                pages.add(int(part))
            except ValueError as exc:
                raise PdfOperationError(f"Invalid page number: {part}") from exc

    return sorted(pages)


def get_pdf_page_count(pdf_path: os.PathLike[str] | str) -> int:
    """Return the total number of pages in the supplied PDF."""
    pdf_path = Path(pdf_path)
    if not pdf_path.exists():
        raise PdfOperationError(f"PDF not found: {pdf_path}")
    try:
        with fitz.open(pdf_path) as doc:
            return len(doc)
    except Exception as exc:
        raise PdfOperationError(f"Unable to read PDF: {pdf_path}") from exc


def merge_pdfs(pdf_paths: Sequence[os.PathLike[str] | str], output_dir: os.PathLike[str] | str,
               prefix: str = "merged_document") -> str:
    """Merge the given PDFs into a single document and return the output filepath."""
    paths = [Path(p) for p in pdf_paths if p]
    if not paths:
        raise PdfOperationError("No PDF files provided for merge")

    for path in paths:
        if not path.exists():
            raise PdfOperationError(f"PDF not found: {path}")

    output_directory = _ensure_output_dir(output_dir)
    timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
    output_path = output_directory / f"{prefix}_{timestamp}.pdf"

    writer = PyPDF2.PdfWriter()
    try:
        for path in paths:
            with path.open("rb") as handle:
                reader = PyPDF2.PdfReader(handle)
                for page in reader.pages:
                    writer.add_page(page)
        with output_path.open("wb") as target:
            writer.write(target)
    except Exception as exc:
        if output_path.exists():
            output_path.unlink(missing_ok=True)
        raise PdfOperationError(f"Failed to merge PDFs: {exc}") from exc

    return str(output_path)


def extract_pages(pdf_path: os.PathLike[str] | str, pages: Iterable[int],
                  output_dir: os.PathLike[str] | str) -> str:
    """Extract specific pages to a new PDF and return the output filepath."""
    pdf_path = Path(pdf_path)
    page_indices = sorted({int(p) for p in pages})
    if not page_indices:
        raise PdfOperationError("No pages to extract")

    if not pdf_path.exists():
        raise PdfOperationError(f"PDF not found: {pdf_path}")

    output_directory = _ensure_output_dir(output_dir)
    if len(page_indices) == 1:
        suffix = f"page_{page_indices[0]}"
    else:
        suffix = f"pages_{page_indices[0]}-{page_indices[-1]}"

    output_path = output_directory / f"extracted_{suffix}.pdf"

    try:
        with fitz.open(pdf_path) as source:
            page_count = len(source)
            for page in page_indices:
                if page < 1 or page > page_count:
                    raise PdfOperationError(
                        f"Page {page} out of range. Document has {page_count} pages."
                    )
            with fitz.open() as target:
                for page in page_indices:
                    target.insert_pdf(source, from_page=page - 1, to_page=page - 1)
                target.save(output_path)
    except PdfOperationError:
        raise
    except Exception as exc:
        if output_path.exists():
            output_path.unlink(missing_ok=True)
        raise PdfOperationError(f"Failed to extract pages: {exc}") from exc

    return str(output_path)


def split_pdf(pdf_path: os.PathLike[str] | str, output_dir: os.PathLike[str] | str) -> str:
    """Split the PDF into individual pages and return the output directory."""
    pdf_path = Path(pdf_path)
    if not pdf_path.exists():
        raise PdfOperationError(f"PDF not found: {pdf_path}")

    try:
        with fitz.open(pdf_path) as source:
            base_name = pdf_path.stem
            output_directory = _ensure_output_dir(Path(output_dir) / f"{base_name}_pages")
            for page_number in range(len(source)):
                output_path = output_directory / f"page_{page_number + 1}.pdf"
                with fitz.open() as target:
                    target.insert_pdf(source, from_page=page_number, to_page=page_number)
                    target.save(output_path)
    except Exception as exc:
        raise PdfOperationError(f"Failed to split PDF: {exc}") from exc

    return str(output_directory)


def rotate_pdf(pdf_path: os.PathLike[str] | str, output_dir: os.PathLike[str] | str,
               rotation: int) -> str:
    """Rotate all pages by the provided angle and return the output filepath."""
    pdf_path = Path(pdf_path)
    if rotation % 90 != 0:
        raise PdfOperationError("Rotation angle must be a multiple of 90 degrees")

    if not pdf_path.exists():
        raise PdfOperationError(f"PDF not found: {pdf_path}")

    rotation = rotation % 360
    output_directory = _ensure_output_dir(output_dir)
    output_path = output_directory / f"{pdf_path.stem}_rotated.pdf"

    try:
        with fitz.open(pdf_path) as source:
            for page in source:
                page.set_rotation(rotation)
            source.save(output_path)
    except Exception as exc:
        if output_path.exists():
            output_path.unlink(missing_ok=True)
        raise PdfOperationError(f"Failed to rotate PDF: {exc}") from exc

    return str(output_path)


def compress_pdf(pdf_path: os.PathLike[str] | str, output_dir: os.PathLike[str] | str) -> tuple[str, float]:
    """Compress the PDF and return the output path plus percentage reduction."""
    pdf_path = Path(pdf_path)
    if not pdf_path.exists():
        raise PdfOperationError(f"PDF not found: {pdf_path}")

    output_directory = _ensure_output_dir(output_dir)
    output_path = output_directory / f"{pdf_path.stem}_compressed.pdf"

    try:
        original_size = pdf_path.stat().st_size
        with fitz.open(pdf_path) as source:
            source.save(
                output_path,
                garbage=4,
                deflate=True,
                clean=True,
                linear=True,
            )
        compressed_size = output_path.stat().st_size
        ratio = 0.0 if original_size == 0 else (1 - compressed_size / original_size) * 100
    except Exception as exc:
        if output_path.exists():
            output_path.unlink(missing_ok=True)
        raise PdfOperationError(f"Failed to compress PDF: {exc}") from exc

    return str(output_path), ratio


def convert_pdf(pdf_path: os.PathLike[str] | str, output_dir: os.PathLike[str] | str,
                output_format: str) -> str:
    """Convert the PDF into the requested format (word, excel, powerpoint, text)."""
    pdf_path = Path(pdf_path)
    if not pdf_path.exists():
        raise PdfOperationError(f"PDF not found: {pdf_path}")

    output_directory = _ensure_output_dir(output_dir)
    fmt = output_format.lower()

    try:
        if fmt == "word":
            output_path = output_directory / f"{pdf_path.stem}.docx"
            cv = Converter(str(pdf_path))
            try:
                cv.convert(str(output_path))
            finally:
                cv.close()
        elif fmt == "excel":
            output_path = output_directory / f"{pdf_path.stem}.xlsx"
            try:
                tables = tabula.read_pdf(str(pdf_path), pages='all')
            except JavaNotFoundError as exc:
                raise PdfOperationError(
                    "Java is required for Excel conversion but was not found. "
                    "Install Java and ensure it is on your PATH."
                ) from exc
            if not tables:
                raise PdfOperationError("No tabular data detected in the PDF")
            with pd.ExcelWriter(output_path) as writer:
                for index, table in enumerate(tables, start=1):
                    table.to_excel(writer, sheet_name=f"Sheet_{index}", index=False)
        elif fmt == "powerpoint":
            output_path = output_directory / f"{pdf_path.stem}.pptx"
            presentation = Presentation()
            with pdfplumber.open(pdf_path) as pdf:
                for page in pdf.pages:
                    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
                    text_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(5))
                    text_frame = text_box.text_frame
                    text = page.extract_text() or ""
                    text_frame.clear()
                    text_frame.text = text
            presentation.save(output_path)
        elif fmt == "text":
            output_path = output_directory / f"{pdf_path.stem}.txt"
            with pdfplumber.open(pdf_path) as pdf, output_path.open("w", encoding="utf-8") as handle:
                for page in pdf.pages:
                    text = page.extract_text() or ""
                    handle.write(text)
                    handle.write("\n\n")
        else:
            raise PdfOperationError(f"Unsupported output format: {output_format}")
    except PdfOperationError:
        raise
    except Exception as exc:
        if 'java' in str(exc).lower():
            raise PdfOperationError(
                "Java is required for Excel conversion but was not found. "
                "Install Java and ensure it is on your PATH."
            ) from exc
        raise PdfOperationError(f"Failed to convert PDF: {exc}") from exc

    return str(output_path)
