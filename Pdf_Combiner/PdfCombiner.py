import os
import json
from PyQt6.QtWidgets import (QApplication, QMainWindow, QPushButton, QFileDialog, QVBoxLayout, 
                            QWidget, QLineEdit, QLabel, QGridLayout, QGroupBox, QStatusBar, 
                            QHBoxLayout, QFrame, QSplitter, QListWidget, QListWidgetItem,
                            QComboBox)  # Add QComboBox import
from PyQt6.QtCore import Qt, QSize
from PyQt6.QtGui import QIcon, QFont, QPixmap, QImage 
import PyPDF2
import fitz 
import datetime
from pdf2docx import Converter
import tabula
import pdfplumber
from pptx import Presentation
from pptx.util import Inches
import pandas as pd
from tabula.errors import JavaNotFoundError


# File persistence functions
def load_last_dir():
    """Load the last used directory from JSON file. Returns current working directory if not found."""
    try:
        with open('last_dir.json', 'r') as f:
            return json.load(f)
    except FileNotFoundError:
        return os.getcwd()


def save_last_dir(dir):
    """Save the last used directory to JSON file for persistence."""
    with open('last_dir.json', 'w') as f:
        json.dump(dir, f)


# PDF core operations
def combine_pdfs(pdf_list, output):
    """Combine multiple PDF files into a single PDF document."""
    try:
        pdf_writer = PyPDF2.PdfWriter()
        for pdf in pdf_list:
            pdf_reader = PyPDF2.PdfReader(pdf)
            for page_num in range(len(pdf_reader.pages)):
                page = pdf_reader.pages[page_num]
                pdf_writer.add_page(page)
        with open(output, 'wb') as out:
            pdf_writer.write(out)
        print("PDFs combined successfully.")
    except Exception as e:
        print("Error combining PDFs:", str(e))


def parse_page_ranges(page_string):
    """
    Parse a page range string into a list of page numbers.
    Accepts formats like "1,3,5-10" and returns a sorted list of unique page numbers.
    Raises ValueError for invalid formats.
    """
    pages = []
    parts = page_string.split(',')
    
    for part in parts:
        part = part.strip()
        if '-' in part:
            # Handle range like "5-10"
            try:
                start, end = part.split('-')
                start = int(start.strip())
                end = int(end.strip())
                if start <= end:
                    pages.extend(range(start, end + 1))
                else:
                    raise ValueError(f"Invalid range: {part} (start > end)")
            except ValueError as e:
                if "invalid literal for int" in str(e):
                    raise ValueError(f"Invalid page range format: {part}")
                raise
        else:
            # Handle single page
            try:
                pages.append(int(part))
            except ValueError:
                raise ValueError(f"Invalid page number: {part}")
    
    return sorted(set(pages))  # Remove duplicates and sort


def extrack_pages(pdf_file, pages, output_dir):
    """
    Extract specific pages from a PDF file and save to a new document using PyMuPDF.
    Returns the output filename on success, None on failure.
    """
    try:
        # Create a filename that includes the page range
        if len(pages) == 1:
            filename = f'extracted_page_{pages[0]}.pdf'
        else:
            # Modified naming format: show first and last page only
            filename = f"extracted_pages_{min(pages)}-{max(pages)}.pdf"
        
        output_filename = os.path.join(output_dir, filename)
        
        # Open the PDF with PyMuPDF
        doc = fitz.open(pdf_file)
        output_doc = fitz.open()
        
        for page_num in pages:
            if 1 <= page_num <= len(doc):
                page = doc.load_page(page_num - 1)  # 0-based indexing
                output_doc.insert_pdf(doc, from_page=page_num-1, to_page=page_num-1)
            else:
                print(f"Warning: Page {page_num} out of range, skipping")
        
        # Save the new document
        output_doc.save(output_filename)
        output_doc.close()
        doc.close()
        
        print(f"Pages extracted successfully to {filename}")
        return output_filename
    except Exception as e:
        print("Error extracting pages:", str(e))
        return None


# PDF manipulation functions
def merge_range(pdf_file, start, end, output_dir):
    """Merge a specific range of pages from a PDF file."""
    try:
        pdf_reader = PyPDF2.PdfReader(pdf_file)
        pdf_writer = PyPDF2.PdfWriter()
        for page_num in range(start, end + 1):
            pdf_writer.add_page(pdf_reader.pages[page_num - 1])
        output_filename = os.path.join(output_dir, 'merged_range.pdf')
        with open(output_filename, 'wb') as out:
            pdf_writer.write(out)
        print("Pages merged successfully.")
    except Exception as e:
        print("Error merging pages:", str(e))


def split_pdf(pdf_file, output_dir):
    """Split a PDF file into individual pages using PyMuPDF."""
    try:
        # Open the PDF with PyMuPDF
        doc = fitz.open(pdf_file)
        num_pages = len(doc)
        
        # Create output directory for pages if it doesn't exist
        pdf_basename = os.path.splitext(os.path.basename(pdf_file))[0]
        output_subdir = os.path.join(output_dir, f"{pdf_basename}_pages")
        os.makedirs(output_subdir, exist_ok=True)
        
        # Extract each page to a separate file
        for page_num in range(num_pages):
            # Create a new PDF with just this page
            output_doc = fitz.open()
            output_doc.insert_pdf(doc, from_page=page_num, to_page=page_num)
            
            # Save the page
            output_filename = os.path.join(output_subdir, f'page_{page_num + 1}.pdf')
            output_doc.save(output_filename)
            output_doc.close()
        
        doc.close()
        print(f"PDF split successfully into {num_pages} pages in folder: {os.path.basename(output_subdir)}")
        return output_subdir
    except Exception as e:
        print("Error splitting PDF:", str(e))
        return None


def rotate_pdf(pdf_file, output_dir, rotation):
    """Rotate all pages in a PDF file by the specified angle using PyMuPDF."""
    try:
        # Create output filename
        original_filename = os.path.basename(pdf_file)
        filename_without_ext, extension = os.path.splitext(original_filename)
        output_filename = os.path.join(output_dir, f'{filename_without_ext}_rotated{extension}')
        
        # Open with PyMuPDF
        doc = fitz.open(pdf_file)
        
        # Apply rotation to each page
        for page_num in range(len(doc)):
            page = doc[page_num]
            page.set_rotation(rotation)
        
        # Save the modified document
        doc.save(output_filename)
        doc.close()
        
        print(f"PDF rotated successfully and saved to {os.path.basename(output_filename)}")
        return output_filename
    except Exception as e:
        print("Error rotating PDF:", str(e))
        return None


def compress_pdf(pdf_file, output_dir):
    """Compress a PDF file to reduce its file size using PyMuPDF (fitz)."""
    try:
        # Create output filename with original name plus "_compressed"
        original_filename = os.path.basename(pdf_file)
        filename_without_ext, extension = os.path.splitext(original_filename)
        output_filename = os.path.join(output_dir, f'{filename_without_ext}_compressed{extension}')
        
        # Open with PyMuPDF
        doc = fitz.open(pdf_file)
        
        # Save with compression options
        doc.save(output_filename, 
                 garbage=4,        # Clean up unused objects
                 deflate=True,     # Use deflate compression
                 clean=True,       # Clean redundancies
                 linear=True)      # Optimize for web
        doc.close()
        
        # Compare file sizes to report compression ratio
        original_size = os.path.getsize(pdf_file)
        compressed_size = os.path.getsize(output_filename)
        ratio = (1 - compressed_size / original_size) * 100
        
        print(f"PDF compressed successfully. Reduced by {ratio:.1f}%")
        return output_filename, ratio
    except Exception as e:
        print("Error compressing PDF:", str(e))
        return None, 0


def convert_pdf(pdf_file, output_dir, output_format):
    """
    Convert PDF to other formats using specialized libraries:
    - pdf2docx for Word conversion
    - tabula-py for Excel conversion (requires Java)
    - python-pptx for PowerPoint conversion
    - pdfplumber for text extraction
    """
    try:
        output_filename = ''
        if output_format == 'word':
            # Convert to Word using pdf2docx
            output_filename = os.path.join(output_dir, 'converted.docx')
            cv = Converter(pdf_file)
            cv.convert(output_filename)
            cv.close()
            
        elif output_format == 'excel':
            # Convert to Excel using tabula-py
            try:
                output_filename = os.path.join(output_dir, 'converted.xlsx')
                # Try to read all tables from the PDF
                tables = tabula.read_pdf(pdf_file, pages='all')
                if not tables:
                    raise ValueError("No tables found in PDF")
                    
                with pd.ExcelWriter(output_filename) as writer:
                    for i, table in enumerate(tables):
                        table.to_excel(writer, sheet_name=f'Sheet_{i+1}', index=False)
                        
            except JavaNotFoundError:
                raise RuntimeError(
                    "Java is required for Excel conversion but was not found.\n"
                    "Please install Java and ensure it's in your system PATH:\n"
                    "1. Download Java from https://www.java.com/download/\n"
                    "2. Install Java\n"
                    "3. Restart your application"
                )
                
        elif output_format == 'powerpoint':
            # Convert to PowerPoint using python-pptx
            output_filename = os.path.join(output_dir, 'converted.pptx')
            prs = Presentation()
            
            # Extract text and create slides
            with pdfplumber.open(pdf_file) as pdf:
                for page in pdf.pages:
                    # Add a slide for each page
                    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Use blank layout
                    text_box = slide.shapes.add_textbox(
                        Inches(1), Inches(1), Inches(8), Inches(5)
                    )
                    text_frame = text_box.text_frame
                    text_frame.text = page.extract_text()
            
            prs.save(output_filename)
            
        elif output_format == 'text':
            # Convert to text using pdfplumber
            output_filename = os.path.join(output_dir, 'converted.txt')
            with pdfplumber.open(pdf_file) as pdf:
                with open(output_filename, 'w', encoding='utf-8') as txt_file:
                    for page in pdf.pages:
                        txt_file.write(page.extract_text())
                        txt_file.write('\n\n')  # Add page breaks
        else:
            raise ValueError("Invalid output format")
            
        print(f"PDF converted successfully to {output_format}")
        return output_filename
        
    except Exception as e:
        error_msg = str(e)
        if "java" in error_msg.lower():
            error_msg = (
                "Java is required for Excel conversion but was not found.\n"
                "Please install Java and ensure it's in your system PATH:\n"
                "1. Download Java from https://www.java.com/en/download/\n"
                "2. Install Java\n" 
                "3. Restart your application"
            )
        print(f"Error converting PDF: {error_msg}")
        raise


def preview_all_pages_pdf(pdf_file):
    try:
        pdf_reader = PyPDF2.PdfReader(pdf_file)
        for page_num in range(len(pdf_reader.pages)):
            page = pdf_reader.pages[page_num]
            print(page.extract_text())
    except Exception as e:
        print("Error previewing PDF:", str(e))


def load_last_dir():
    try:
        with open('last_dir.json', 'r') as f:
            return json.load(f)
    except FileNotFoundError:
        return os.getcwd()


def save_last_dir(dir):
    with open('last_dir.json', 'w') as f:
        json.dump(dir, f)


def drop_pdf_files(event, pdf_files):
    if event.mimeData().hasUrls():
        event.setDropAction(Qt.DropAction.ActionCopyAction)
        event.accept()
        for url in event.mimeData().urls():
            pdf_files.append(url.toLocalFile())
    else:
        event.ignore()


def drop_position(event):
    if event.mimeData().hasUrls():
        event.setDropAction(Qt.DropAction.ActionCopyAction)
        event.accept()
        position = event.pos()
    else:
        event.ignore()


def drop_leave(event):
    event.accept()


def batch_process(pdf_files, output_dir, process_func):
    for pdf_file in pdf_files:
        process_func(pdf_file, output_dir)


def batch_process_with_args(pdf_files, output_dir, process_func, *args):
    for pdf_file in pdf_files:
        process_func(pdf_file, output_dir, *args)


def batch_process_with_args_and_output(pdf_files, output_dir, process_func, *args):
    for pdf_file in pdf_files:
        process_func(pdf_file, output_dir, *args)


# PDF preview and rendering
def render_page_to_pixmap(pdf_path, page_num=0, size=None):
    """
    Render a PDF page to a QPixmap for preview.
    Args:
        pdf_path: Path to PDF file
        page_num: Page number to render (0-based)
        size: Optional target size for the pixmap
    Returns: QPixmap object or None on failure
    """
    try:
        # Open the PDF with PyMuPDF (fitz)
        doc = fitz.open(pdf_path)
        if page_num >= len(doc):
            return None
            
        page = doc[page_num]
        
        # Set a reasonable zoom factor
        zoom = 2.0  # Increase for higher quality
        mat = fitz.Matrix(zoom, zoom)
        
        # Render page to pixmap
        pix = page.get_pixmap(matrix=mat)
        
        # Convert to QImage
        img_format = QImage.Format.Format_RGB888  # Note: Use Format.Format_RGB888, not Format_RGB888
        img = QImage(pix.samples, pix.width, pix.height, pix.stride, img_format)
        
        # Convert to QPixmap
        pixmap = QPixmap.fromImage(img)
        
        # Scale if size is provided
        if size:
            pixmap = pixmap.scaled(size, Qt.AspectRatioMode.KeepAspectRatio, 
                                 Qt.TransformationMode.SmoothTransformation)
        
        return pixmap
    except Exception as e:
        print(f"Error rendering PDF: {str(e)}")
        return None


# Main application
def main():
    """
    Main application entry point. Creates and configures the PDF Toolkit GUI.
    Features:
    - PDF file selection and management
    - Page preview and navigation
    - Multiple PDF operations (merge, split, rotate, extract)
    - Drag and drop support
    - Output directory management
    """
    app = QApplication([])
    window = QMainWindow()
    window.setWindowTitle("PDF Toolkit")
    window.setGeometry(100, 100, 900, 700)
    window.setStyleSheet("""
        QMainWindow, QWidget {
            background-color: #0d0936;
            color: white;
        }
        QGroupBox {
            font-weight: bold;
            border: 1px solid #2a2a5a;
            border-radius: 5px;
            margin-top: 10px;
            padding-top: 15px;
            color: white;
        }
        QPushButton {
            background-color: #4a86e8;
            color: white;
            border: none;
            padding: 8px 16px;
            border-radius: 4px;
            font-size: 12px;
        }
        QPushButton:hover {
            background-color: #3a76d8;
        }
        QLineEdit {
            padding: 5px;
            border: 1px solid #2a2a5a;
            border-radius: 3px;
            background-color: #030d22;
            color: white;
        }
        QLabel {
            font-size: 12px;
            color: white;
        }
        QListWidget {
            background-color: #030d22;
            color: white;
            border: 1px solid #2a2a5a;
        }
        QStatusBar {
            background-color: #030d22;
            color: white;
        }
    """)

    # Main layout with splitter
    main_widget = QWidget()
    main_layout = QHBoxLayout(main_widget)
    main_layout.setContentsMargins(20, 20, 20, 20)
    main_layout.setSpacing(15)
    
    splitter = QSplitter(Qt.Orientation.Horizontal)
    
    # Left side widget for controls
    left_widget = QWidget()
    left_layout = QVBoxLayout(left_widget)
    left_layout.setContentsMargins(0, 0, 0, 0)
    left_layout.setSpacing(15)

    # File selection section
    file_group = QGroupBox("File Selection")
    file_layout = QGridLayout()
    
    pdf_files_button = QPushButton("Browse PDF Files")
    pdf_files_button.setMinimumHeight(35)
    
    file_list = QListWidget()
    file_list.setMaximumHeight(100)
    # Enable checkboxes and drag-drop
    file_list.setSelectionMode(QListWidget.SelectionMode.MultiSelection)
    file_list.setDragDropMode(QListWidget.DragDropMode.InternalMove)
    
    # Add Select All button to file order controls
    file_order_layout = QHBoxLayout()
    select_all_button = QPushButton("Select All")
    move_up_button = QPushButton("↑")
    move_up_button.setMaximumWidth(40)
    move_down_button = QPushButton("↓")
    move_down_button.setMaximumWidth(40)
    remove_file_button = QPushButton("Remove")
    
    file_order_layout.addWidget(select_all_button)
    file_order_layout.addWidget(move_up_button)
    file_order_layout.addWidget(move_down_button)
    file_order_layout.addWidget(remove_file_button)
    
    output_dir_button = QPushButton("Set Output Directory")
    output_dir_button.setMinimumHeight(35)
    
    output_label = QLabel("Output directory: Not selected")
    
    file_layout.addWidget(pdf_files_button, 0, 0)
    file_layout.addWidget(file_list, 1, 0, 1, 2)
    file_layout.addLayout(file_order_layout, 2, 0, 1, 2)  # Add the order controls
    file_layout.addWidget(output_dir_button, 3, 0)
    file_layout.addWidget(output_label, 3, 1)
    
    file_group.setLayout(file_layout)
    
    # PDF Operations section
    operations_group = QGroupBox("PDF Operations")
    operations_layout = QGridLayout()
    
    # Extract section (now at position 0)
    extract_label = QLabel("Pages to extract (e.g. 1,3,5):")
    pages_line_edit = QLineEdit()
    extract_button = QPushButton("Extract Pages")
    
    operations_layout.addWidget(extract_label, 0, 0)
    operations_layout.addWidget(pages_line_edit, 0, 1)
    operations_layout.addWidget(extract_button, 0, 2)
    
    # Merge section (now at position 1)
    merge_button = QPushButton("Merge Selected PDFs")
    operations_layout.addWidget(merge_button, 1, 0)
    
    # Split section (now at position 2)
    split_button = QPushButton("Split PDF")
    operations_layout.addWidget(split_button, 2, 0)
    
    # Rotate section (now at position 3)
    rotation_label = QLabel("Rotation angle (90, 180, 270):")
    rotation_line_edit = QLineEdit()
    rotate_button = QPushButton("Rotate PDF")
    
    operations_layout.addWidget(rotation_label, 3, 0)
    operations_layout.addWidget(rotation_line_edit, 3, 1)
    operations_layout.addWidget(rotate_button, 3, 2)

    # Compress section (now at position 4)
    compress_button = QPushButton("Compress PDF")
    operations_layout.addWidget(compress_button, 4, 0)
    
    # Convert section (now at position 5)
    convert_label = QLabel("Convert to format:")
    format_combo = QComboBox()
    format_combo.addItems(['Word', 'Excel', 'PowerPoint', 'Text'])
    convert_button = QPushButton("Convert PDF")
    
    operations_layout.addWidget(convert_label, 5, 0)
    operations_layout.addWidget(format_combo, 5, 1)
    operations_layout.addWidget(convert_button, 5, 2)
    
    operations_group.setLayout(operations_layout)
    
    # Add sections to left layout
    left_layout.addWidget(file_group)
    left_layout.addWidget(operations_group)
    
    # Right side widget for preview
    right_widget = QWidget()
    right_layout = QVBoxLayout(right_widget)
    right_layout.setContentsMargins(0, 0, 0, 0)
    right_layout.setSpacing(15)
    
    # PDF Preview section
    preview_group = QGroupBox("PDF Preview")
    preview_layout = QVBoxLayout()
    
    # Create preview container with fixed size
    preview_container = QWidget()
    preview_container.setMinimumSize(400, 500)
    preview_container.setMaximumHeight(500)
    preview_container_layout = QVBoxLayout(preview_container)
    preview_container_layout.setContentsMargins(0, 0, 0, 0)
    
    preview_label = QLabel("No PDF selected")
    preview_label.setAlignment(Qt.AlignmentFlag.AlignCenter)
    preview_label.setStyleSheet("background-color: #030d22; border: 1px solid #2a2a5a;")
    preview_container_layout.addWidget(preview_label)
    
    # Navigation buttons with fixed width container
    nav_container = QWidget()
    nav_container.setFixedHeight(50)  # Fixed height for navigation
    nav_layout = QHBoxLayout(nav_container)
    nav_layout.setContentsMargins(0, 5, 0, 5)
    
    # Set fixed widths for navigation elements
    prev_button = QPushButton("Previous Page")
    prev_button.setFixedWidth(120)
    
    page_nav_layout = QHBoxLayout()
    page_nav_layout.setSpacing(5)
    page_input = QLineEdit()
    page_input.setFixedWidth(50)
    page_input.setPlaceholderText("Page")
    page_label = QLabel("/")
    total_pages_label = QLabel("0")
    total_pages_label.setFixedWidth(30)
    
    next_button = QPushButton("Next Page")
    next_button.setFixedWidth(120)
    
    # Add widgets to page navigation
    page_nav_layout.addWidget(page_input)
    page_nav_layout.addWidget(page_label)
    page_nav_layout.addWidget(total_pages_label)
    
    # Add all elements to navigation layout
    nav_layout.addWidget(prev_button)
    nav_layout.addStretch()
    nav_layout.addLayout(page_nav_layout)
    nav_layout.addStretch()
    nav_layout.addWidget(next_button)
    
    # Add everything to preview layout
    preview_layout.addWidget(preview_container)
    preview_layout.addWidget(nav_container)
    
    preview_group.setLayout(preview_layout)
    
    # Add preview to right layout
    right_layout.addWidget(preview_group)
    
    # Add widgets to splitter
    splitter.addWidget(left_widget)
    splitter.addWidget(right_widget)
    
    # Set initial sizes
    splitter.setSizes([400, 500])
    
    # Add splitter to main layout
    main_layout.addWidget(splitter)
    
    # Status bar
    status_bar = QStatusBar()
    window.setStatusBar(status_bar)
    status_bar.showMessage("Ready")

    # Variables
    pdf_files = []
    output_dir = load_last_dir()
    current_page = 0
    total_pages = 0
    
    # Event handlers
    def browse_pdf_files():
        """Handle PDF file selection dialog and update the file list."""
        file_dialog = QFileDialog()
        file_dialog.setFileMode(QFileDialog.FileMode.ExistingFiles)
        file_dialog.setNameFilter("PDF Files (*.pdf)")
        if file_dialog.exec():
            new_files = file_dialog.selectedFiles()
            for file in new_files:
                if file not in pdf_files:
                    pdf_files.append(file)
                    item = QListWidgetItem(os.path.basename(file))
                    item.setFlags(item.flags() | Qt.ItemFlag.ItemIsUserCheckable)
                    item.setCheckState(Qt.CheckState.Unchecked)
                    file_list.addItem(item)
            
            if pdf_files and file_list.count() > 0:
                file_list.setCurrentRow(0)  # Select first item
                update_preview()  # Update preview with the selected file
            status_bar.showMessage(f"{len(new_files)} PDF files added")

    def update_preview():
        """Update the PDF preview display with the currently selected file and page."""
        nonlocal current_page, total_pages
        selected_items = file_list.selectedItems()
        if not selected_items or not pdf_files:  # Add check for empty pdf_files
            preview_label.setText("No PDF selected")
            current_page = 0
            total_pages = 0
            page_input.setText("")
            total_pages_label.setText("0")
            return
            
        selected_pdf = pdf_files[file_list.row(selected_items[0])]
        try:
            # Open PDF to get page count
            with fitz.open(selected_pdf) as doc:
                total_pages = len(doc)
                if current_page >= total_pages:
                    current_page = 0
                
                # Update page indicators
                page_input.setText(str(current_page + 1))
                total_pages_label.setText(str(total_pages))
                
                # Render the page
                pixmap = render_page_to_pixmap(selected_pdf, current_page, preview_label.size())
                if (pixmap):
                    preview_label.setPixmap(pixmap)
                else:
                    preview_label.setText("Unable to render page")
        except Exception as e:
            preview_label.setText(f"Error loading PDF: {str(e)}")
            print(f"Preview error: {str(e)}")
            total_pages = 0
            current_page = 0
            page_input.setText("")
            total_pages_label.setText("0")

    def go_to_page():
        nonlocal current_page
        try:
            page_num = int(page_input.text()) - 1  # Convert to 0-based index
            if 0 <= page_num < total_pages:
                current_page = page_num
                update_preview()
            else:
                status_bar.showMessage(f"Page number must be between 1 and {total_pages}")
        except ValueError:
            status_bar.showMessage("Please enter a valid page number")

    def move_file_up():
        current_row = file_list.currentRow()
        if current_row > 0:
            # Move item in the list widget
            current_item = file_list.takeItem(current_row)
            file_list.insertItem(current_row - 1, current_item)
            file_list.setCurrentRow(current_row - 1)
            
            # Also reorder in the pdf_files list
            pdf_files.insert(current_row - 1, pdf_files.pop(current_row))
            status_bar.showMessage(f"Moved file up: {os.path.basename(pdf_files[current_row - 1])}")

    def move_file_down():
        current_row = file_list.currentRow()
        if current_row < file_list.count() - 1:
            # Move item in the list widget
            current_item = file_list.takeItem(current_row)
            file_list.insertItem(current_row + 1, current_item)
            file_list.setCurrentRow(current_row + 1)
            
            # Also reorder in the pdf_files list
            pdf_files.insert(current_row + 1, pdf_files.pop(current_row))
            status_bar.showMessage(f"Moved file down: {os.path.basename(pdf_files[current_row + 1])}")

    def remove_selected_file():
        current_row = file_list.currentRow()
        if current_row >= 0:
            # First remove from UI
            file_list.takeItem(current_row)
            # Then remove from data list
            if current_row < len(pdf_files):
                removed_file = pdf_files.pop(current_row)
                status_bar.showMessage(f"Removed file: {os.path.basename(removed_file)}")
            
            # Update UI state
            if file_list.count() > 0:
                # Select next available item
                new_row = min(current_row, file_list.count() - 1)
                file_list.setCurrentRow(new_row)
            else:
                # Reset UI when no items left
                preview_label.setText("No PDF selected")
                current_page = 0
                total_pages = 0
                page_input.setText("")
                total_pages_label.setText("0")

    def show_next_page():
        nonlocal current_page
        if total_pages > 0 and current_page < total_pages - 1:
            current_page += 1
            update_preview()

    def show_prev_page():
        nonlocal current_page
        if current_page > 0:
            current_page -= 1
            update_preview()

    def browse_output_dir():
        nonlocal output_dir
        dir_path = QFileDialog.getExistingDirectory()
        if dir_path:
            output_dir = dir_path
            save_last_dir(output_dir)
            output_label.setText(f"Output directory: {output_dir}")
            status_bar.showMessage(f"Output directory set to: {output_dir}")

    def combine():
        if not pdf_files:
            status_bar.showMessage("No PDF files selected")
            return
        if not output_dir:
            status_bar.showMessage("No output directory selected")
            return
            
        try:
            pdf_writer = PyPDF2.PdfWriter()
            
            # Add all pages from all selected PDFs to the writer
            for pdf_file in pdf_files:
                with open(pdf_file, 'rb') as f:
                    pdf_reader = PyPDF2.PdfReader(f)
                    for page_num in range(len(pdf_reader.pages)):
                        page = pdf_reader.pages[page_num]
                        pdf_writer.add_page(page)
            
            output_filename = os.path.join(output_dir, 'combined_document.pdf')
            with open(output_filename, 'wb') as out:
                pdf_writer.write(out)
                
            status_bar.showMessage(f"All PDF files combined successfully to {output_filename}")
        except Exception as e:
            status_bar.showMessage(f"Error combining PDFs: {str(e)}")

    def extract():
        if not pdf_files:
            status_bar.showMessage("No PDF files selected")
            return
            
        pages_text = pages_line_edit.text()
        if not pages_text:
            status_bar.showMessage("Please enter pages to extract")
            return
            
        try:
            page_list = parse_page_ranges(pages_text)
            if not page_list:
                status_bar.showMessage("No valid pages specified")
                return
                
            output_file = extrack_pages(pdf_files[0], page_list, output_dir)
            if output_file:
                filename = os.path.basename(output_file)
                status_bar.showMessage(f"Extracted pages {pages_text} to {filename}")
        except Exception as e:
            status_bar.showMessage(f"Error extracting pages: {str(e)}")

    def merge():
        if not pdf_files:
            status_bar.showMessage("No PDF files selected")
            return
        if not output_dir:
            status_bar.showMessage("No output directory selected")
            return
        
        try:
            pdf_writer = PyPDF2.PdfWriter()
            files_to_merge = []
            
            # Get all checked files
            for i in range(file_list.count()):
                item = file_list.item(i)
                if item.checkState() == Qt.CheckState.Checked:
                    files_to_merge.append(pdf_files[i])
            
            if not files_to_merge:
                status_bar.showMessage("No files selected for merging")
                return
            
            # Merge all checked files
            for pdf_file in files_to_merge:
                try:
                    with open(pdf_file, 'rb') as f:
                        pdf_reader = PyPDF2.PdfReader(f)
                        for page in pdf_reader.pages:
                            pdf_writer.add_page(page)
                except Exception as e:
                    status_bar.showMessage(f"Error processing {os.path.basename(pdf_file)}: {str(e)}")
                    return
            
            timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
            output_filename = os.path.join(output_dir, f'merged_document_{timestamp}.pdf')
            
            with open(output_filename, 'wb') as out:
                pdf_writer.write(out)
            
            status_bar.showMessage(f"Successfully merged {len(files_to_merge)} PDF files to {os.path.basename(output_filename)}")
        except Exception as e:
            status_bar.showMessage(f"Error merging PDFs: {str(e)}")

    def split():
        if not pdf_files:
            status_bar.showMessage("Error: No PDF files selected")
            return
            
        try:
            split_pdf(pdf_files[0], output_dir)
            status_bar.showMessage("PDF split successfully")
        except Exception as e:
            status_bar.showMessage(f"Error splitting PDF: {str(e)}")

    def rotate():
        if not pdf_files:
            status_bar.showMessage("Error: No PDF files selected")
            return
            
        try:
            rotation = int(rotation_line_edit.text())
            rotate_pdf(pdf_files[0], output_dir, rotation)
            status_bar.showMessage("PDF rotated successfully")
        except Exception as e:
            status_bar.showMessage(f"Error rotating PDF: {str(e)}")
    
    def compress():
        if not pdf_files:
            status_bar.showMessage("Error: No PDF files selected")
            return
            
        try:
            compress_pdf(pdf_files[0], output_dir)
            status_bar.showMessage("PDF compressed successfully")
        except Exception as e:
            status_bar.showMessage(f"Error compressing PDF: {str(e)}")

    def convert():
        if not pdf_files:
            status_bar.showMessage("Error: No PDF files selected")
            return
            
        output_format = format_combo.currentText().lower()
        try:
            convert_pdf(pdf_files[0], output_dir, output_format)
            status_bar.showMessage(f"PDF converted to {output_format} successfully")
        except Exception as e:
            status_bar.showMessage(f"Error converting PDF: {str(e)}")

    # Add select all functionality
    def toggle_select_all():
        all_checked = True
        for i in range(file_list.count()):
            item = file_list.item(i)
            if item.checkState() == Qt.CheckState.Unchecked:
                all_checked = False
                break
        
        new_state = Qt.CheckState.Unchecked if all_checked else Qt.CheckState.Checked
        for i in range(file_list.count()):
            file_list.item(i).setCheckState(new_state)
        
        status_bar.showMessage("All files " + ("deselected" if all_checked else "selected"))

    # Connect signals
    pdf_files_button.clicked.connect(browse_pdf_files)
    output_dir_button.clicked.connect(browse_output_dir)
    extract_button.clicked.connect(extract)
    merge_button.clicked.connect(merge)
    split_button.clicked.connect(split)
    rotate_button.clicked.connect(rotate)
    prev_button.clicked.connect(show_prev_page)
    next_button.clicked.connect(show_next_page)
    
    # Connect new file order buttons
    move_up_button.clicked.connect(move_file_up)
    move_down_button.clicked.connect(move_file_down)
    remove_file_button.clicked.connect(remove_selected_file)
    select_all_button.clicked.connect(toggle_select_all)
    
    # Connect page input
    page_input.returnPressed.connect(go_to_page)

    compress_button.clicked.connect(compress)
    convert_button.clicked.connect(convert)
    
    # Update preview when a file is clicked in the list
    file_list.itemClicked.connect(update_preview)
    file_list.currentItemChanged.connect(update_preview)
    
    # Handle preview resize events
    def resize_event():
        if pdf_files:
            update_preview()
    
    preview_label.resizeEvent = lambda event: resize_event()
    
    # Event handlers for drag and drop reordering
    def handle_item_moved():
        """Synchronize the internal PDF files list with the UI list after drag-drop reordering."""
        # Update the pdf_files list to match the new order in the list widget
        nonlocal pdf_files
        new_pdf_files = []
        for i in range(file_list.count()):
            item_text = file_list.item(i).text()
            # Find the file path that matches this filename
            for file_path in pdf_files:
                if os.path.basename(file_path) == item_text:
                    new_pdf_files.append(file_path)
                    break
        pdf_files = new_pdf_files
        status_bar.showMessage("Files reordered")

    # Connect the model's data changed signal to handle reordering
    file_list.model().rowsMoved.connect(handle_item_moved)
    
    # Set output directory label if it exists
    if output_dir:
        output_label.setText(f"Output directory: {output_dir}")

    window.setCentralWidget(main_widget)
    window.show()
    app.exec()


# Entry point
if __name__ == '__main__':
    main()